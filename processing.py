
"""
data processing utilities
"""

def get_binary_contiguous_indices(array, partial=False):
    """
    identify indices of array where binary values are True and contiguous
    e.g. array = np.array([False, False, True, True, True, False, True, True, False, True, True])
        returns
        [(2, 3, 4), (6, 7)]
        if partial=False (9, 10) is not part of the result because it is not clear that contiguous block has ended
        otherwise (9, 10) is part of the result
    """
    import numpy as np

    # validate array and convert to int
    assert array.dtype == np.bool
    assert array.ndim == 1
    array = array.astype(np.int)

    # identify indices of positive and negative transitions
    plus = np.where(np.diff(array) == 1)[0] + 1
    minus = np.where(np.diff(array) == -1)[0]

    # case where array is True at both ends
    if (plus.size == minus.size) and (array[0] == 1):
        assert array[-1] == 1
        assert minus[0] < plus[0]
        assert plus[-1] > minus[-1]

        if partial == True:
            plus = np.hstack((0, plus))
            minus = np.hstack((minus, array.size - 1))
        else:
            minus = minus[1:]
            plus = plus[:-1]

    # case where array is True at one end but not the other
    # (does not handle the partial case)
    if plus.size != minus.size:
        assert np.logical_xor(array[0] == 1, array[-1] == 1)

        # array is True at the left end
        if array[0] == 1:
            minus = minus[1:]

        # array is True at the right end
        if array[-1] == 1:
            plus = plus[:-1]
        assert plus.size == minus.size

    # combine and return plus and minus indices
    return [(a, b) for a, b in zip(plus, minus)]

def int_to_ordinal(x):
    """
    converts 1 to '1st', 2 to '2nd', ...
    copied from https://stackoverflow.com/questions/9647202/ordinal-numbers-replacement
    """
    return '%d%s' % (x,'tsnrhtdd'[(x//10%10!=1)*(x%10<4)*x%10::4])

def webm_to_gps(lon, lat):
    """
    convert web-mercator coordinates (EPSG-3857) to GPS coordinates (EPSG-4326)
    returns lon, lat in Web Mercator coords
    """
    from pyproj import Transformer
    x = Transformer.from_crs(crs_from=3857, crs_to=4326, always_xy=True).transform(lon, lat)
    return x[0], x[1]

def gps_to_webm(lon, lat):
    """
    convert GPS coordinates (EPSG-4326) to web-mercator coordinates (EPSG-3857)
    returns lon, lat in GPS coords
    """
    import numpy as np
    from pyproj import Transformer

    # validate
    c0 = isinstance(lon, (int, float))
    c1 = isinstance(lat, (int, float))
    c2 = isinstance(lon, np.ndarray)
    c3 = isinstance(lat, np.ndarray)
    assert (c0 and c1) or (c2 and c3 and (lon.size == lat.size))

    # from arcgis.gis import GIS
    # from arcgis.geometry import project
    # gis = GIS()
    # xy = project([[lon, lat]], in_sr=4326, out_sr=3857)[0]
    # return xy['x'], xy['y']

    # convert and return
    xy = Transformer.from_crs(crs_from=4326, crs_to=3857, always_xy=True).transform(lon, lat)
    return xy[0], xy[1]

def reverse_geocode(lon, lat):
    """
    reverse geocode based on lon / lat, return '<city>, <state>' as a str
    - both arcgis and reverse_geocoder libraries interfere with datalake connections
    """
    import numpy as np
    from arcgis.gis import GIS
    from arcgis import geocoding

    # validate
    c0 = isinstance(lon, (int, float))
    c1 = isinstance(lat, (int, float))
    c2 = isinstance(lon, np.ndarray)
    c3 = isinstance(lat, np.ndarray)
    assert (c0 and c1) or (c2 and c3 and (lon.size == lat.size))

    # arcgis reverse geocode
    gis = GIS()
    if (c0 and c1):
        loc = geocoding.reverse_geocode([lon, lat])
        return f"""{loc['address']['City']}, {loc['address']['RegionAbbr']}"""
    loc = [geocoding.reverse_geocode([lx, ly]) for lx, ly in zip(lon, lat)]
    return np.array([' '.join((x['address']['City'], x['address']['RegionAbbr'])) for x in loc])

    # use reverse_geocoder library
    # import reverse_geocoder
    # rg = reverse_geocoder.RGeocoder()
    # locations = rg.query([(a, b) for a, b in zip(lat, lon)])

def get_bounds_of_data_within_interval(data, x=0.95):
    """
    return bounds of an array 'data' where at least x percent of the data are contained
    - bounds will be items from data for which the interval is inclusive

    e.g. get_bounds_of_data_within_interval([1, 2, 3, 4, 5, 6, 7, 8, 9, 10], x=0.50) -> (3, 7)
         get_bounds_of_data_within_interval([1, 2, 3, 4, 5, 6, 7, 800, 900, 1000], x=0.50) -> (4, 800)
         get_bounds_of_data_within_interval([-1, -2, -3, 4, 5, 6, 7, 8, 9, 10], x=0.50) -> (-1, 7)
    """
    import numpy as np

    data = np.array(data).flatten()
    data = np.sort(data)
    remove = int(data.size * (1 - x))
    assert (data.size - remove) / data.size >= x
    assert (data.size - remove - 1) / data.size < x
    if remove % 2 == 0:     # remove same number of data points from each side of data
        remove = int(remove / 2)
        return data[remove], data[-remove - 1]
    else:                   # decide whether to remove an additional data point from left or right side - maximize interval as the constraint
        remove = int(remove / 2)
        left_diff = np.diff(data[remove : remove + 2])
        right_diff = np.diff(data[-remove - 2 : -remove])
        if left_diff < right_diff:
            return data[remove + 1], data[-remove - 1]
        else:
            return data[remove], data[-remove - 2]

def convert_all_jpg_2_pdf(datadir):
    """
    convert all .jpg files in datadir to .pdf files with matching filenames
    """

    import img2pdf
    import os
    from glob import glob

    jpgs = glob(os.path.join(datadir, '*.jpg'))
    for jpg in jpgs:
        fn = jpg[:-4] + '.pdf'
        with open(fn, 'wb') as fid:
            fid.write(img2pdf.convert(jpg))

def create_filename(fn):
    """
    create a sane / safe filename from fn
    """
    keepcharacters = (' ','.','_')
    return "".join(c for c in fn if c.isalnum() or c in keepcharacters).rstrip()

def get_aliased_freq(f, fs):
    """
    return aliased frequency of f sampled at fs
    """
    import numpy as np

    fn = fs / 2
    if np.int(f / fn) % 2 == 0:
        return f % fn
    else:
        return fn - (f % fn)

def matlab2datetime(mat_dnums):
    """
    convert a 1d numpy array of float matlab datenums (e.g., 735976.001453)
    to a 1d numpy array of datetime objects
    """
    from datetime import datetime, timedelta
    import numpy as np

    day = np.array([datetime.fromordinal(int(x)) for x in mat_dnums])
    dayfrac = np.array([timedelta(days=(x % 1)) - timedelta(days=366) for x in mat_dnums])
    return day + dayfrac

def numpy_datetime64_to_datetime(dt):
    """
    convert an array of numpy datetime64 objects to an array of datetime objects
    """
    import numpy as np
    from datetime import datetime

    assert isinstance(dt, np.ndarray)
    assert dt.dtype.type == np.datetime64

    tref = np.datetime64('1970-01-01T00:00:00Z')
    return np.array([datetime.utcfromtimestamp((x - tref) / np.timedelta64(1, 's')) for x in dt])

def get_folder_size(start_path):
    """
    get size of folder in bytes, copied from
    https://stackoverflow.com/questions/1392413/calculating-a-directorys-size-using-python
    """
    import os

    total_size = 0
    for dirpath, dirnames, filenames in os.walk(start_path):
        for f in filenames:
            fp = os.path.join(dirpath, f)
            # skip if it is symbolic link
            if not os.path.islink(fp):
                total_size += os.path.getsize(fp)

    return total_size

def linspace(start, stop, num):
    """
    exactly the same as numpy.linspace if start and stop are both not datetime objects
    if start and stop are datetime objects:
        * return an array of datetime objects with 'num' elements bounded by start and stop
    """
    import numpy as np
    from datetime import datetime, timedelta

    # use np.linspace in case start and stop are both not datetime objects
    if (not isinstance(start, datetime)) or (not isinstance(start, datetime)):
        return np.linspace(start, stop, num)

    # convert start and stop to seconds since a reference time
    tref = datetime.strptime('1 Jan 1900', '%d %b %Y')
    start = (start - tref).total_seconds()
    stop = (stop - tref).total_seconds()

    # create the time array and convert back to datetime objects
    time = np.linspace(start, stop, num)
    return np.array([tref + timedelta(seconds=x) for x in time])

def arange(start, stop, step):
    """
    exactly the same as numpy.arange is start, stop, and step are all not datetime / timedelta objects
    if start and stop are datetime objects, and step is a timedelta object:
        * return an array of datetime objects between start (inclusive) and stop (maybe inclusive) at step intervals
    (numpy.arange does work with datetime / timedelta objects, but it converts output array to numpy native datetime objects, not datetimes...)
    """

    import numpy as np
    from datetime import datetime, timedelta

    # use np.linspace in case start and stop are both not datetime objects
    if (not isinstance(start, datetime)) or (not isinstance(start, datetime)) or (not isinstance(step, timedelta)):
        return np.arange(start, stop, step)

    # convert start and stop to seconds since a reference time
    tref = datetime.strptime('1 Jan 1900', '%d %b %Y')
    start = (start - tref).total_seconds()
    stop = (stop - tref).total_seconds()

    # create the time array and convert back to datetime objects
    time = np.arange(start, stop, step.total_seconds())
    return np.array([tref + timedelta(seconds=x) for x in time])

def is_datetime_week_number_begin(dt):
    """
    determine if the datetime object dt represents the beginning of an isocalendar week number
    """

    import numpy as np
    from datetime import datetime, timedelta

    week = arange(dt, dt + timedelta(days=7), timedelta(days=1))
    workweek = [x.isocalendar()[1] for x in week]
    if np.unique(workweek).size == 1:
        return True
    else:
        print('{} does not represent the beginning of an isocalendar workweek\ntry {} instead'.format(
            dt.strftime(r'%d %b %Y'), week[np.where(np.diff(workweek))[0][0] + 1].strftime(r'%d %b %Y')))
        return False

def is_datetime_week_number_end(dt):
    """
    determine if the datetime object dt represents the end of an isocalendar week number
    """

    import numpy as np
    from datetime import datetime, timedelta

    week = arange(dt - timedelta(days=6), dt + timedelta(days=1), timedelta(days=1))
    workweek = [x.isocalendar()[1] for x in week]
    if np.unique(workweek).size == 1:
        return True
    else:
        print('{} does not represent the end of an isocalendar workweek\ntry {} instead'.format(
            dt.strftime(r'%d %b %Y'), week[np.where(np.diff(workweek) == 1)[0][0]].strftime(r'%d %b %Y')))
        return False

def datetime_to_week_number(dt):
    """
    convert a datetime object to an isocalendar week number in the format '2018-33'
    information about the day of the week is lost in this conversion
    """
    dt = dt.isocalendar()
    return '{}-{:02d}'.format(dt[0], dt[1])

def week_number_to_datetime(week, day_of_week=1):
    """
    convert an isocalendar week number in the format '2018-17' to a datetime object
    the day of the week for each week number must be provided as input, or the first day is always assumed
    """
    from datetime import datetime

    week += '-{}'.format(day_of_week)
    return datetime.strptime(week, r'%Y-%W-%w')

def loadmat(fname):
    """
    this function should be called instead of directly calling scipy.io.loadmat
    as it solves the problem of not properly recovering python dictionaries
    """

    import scipy.io as spio

    def check_keys(dict):
        """
        checks if entries in dictionary are mat-objects. If yes
        todict is called to change them to nested dictionaries
        """
        for key in dict:
            if isinstance(dict[key], spio.matlab.mio5_params.mat_struct):
                dict[key] = todict(dict[key])
        return dict

    def todict(matobj):
        """
        A recursive function which constructs from matobjects nested dictionaries
        """
        dict = {}
        for strg in matobj._fieldnames:
            elem = matobj.__dict__[strg]
            if isinstance(elem, spio.matlab.mio5_params.mat_struct):
                dict[strg] = todict(elem)
            else:
                dict[strg] = elem
        return dict

    data = spio.loadmat(fname, struct_as_record=False, squeeze_me=True)
    return check_keys(data)

def run_notch_filter_example():
    """
    run a notch filter example
    """

    import matplotlib.pyplot as plt
    import numpy as np
    import pyrb
    from scipy import signal
    from pyrb.mpl import open_figure, format_axes, largefonts
    plt.style.use('bmh')

    # define a sampling rate, fs, and N data points
    fs = 6000
    N = 1e5

    # calculate a time array based on fs and N
    dt = 1 / fs
    time = np.arange(0, N*dt, dt)

    # define y(time) data to includes freqs at mags, plus some baseline noise
    mags = [1, 2, 4, 2, 5, 3, 1]
    freqs = [250, 1200, 1917, 711, 2356, 2100, 8209]
    y = 0
    for mag, freq in zip(mags, freqs):
        y += mag * np.sin(2 * np.pi * freq * time)
    y += np.random.normal(0, 1, y.size)

    # calculate the psd of y data
    freq, psd = signal.welch(y, fs=fs, nperseg=512)

    # update freqs for aliasing, as any freq greater than fs/2 will alias to some other freq less than fs/2
    freqs = [get_aliased_freq(x, fs) for x in freqs]

    # select a random 'freqs' to filter, mapped to 0 to 1 scale where fs/2 maps to 1
    wf = np.random.choice(freqs) / (fs/2)

    # prepare the 0 to 1 mapped wp (pass-band) and ws (stop-band) edge frequencies
    wd = 25 / (fs/2)
    ws = [wf - wd, wf + wd]
    wp = [wf - 2 * wd, wf + 2 * wd]
    gpass, gstop = 3, 40

    # create the bandstop filter
    N, Wn = signal.cheb2ord(wp=wp, ws=ws, gpass=gpass, gstop=gstop)
    b, a = signal.iirfilter(N=N, Wn=Wn, rp=gpass, rs=gstop, btype='bandstop', ftype='cheby2')

    # apply the filter to y, get the psd of the filtered signal
    yf = signal.lfilter(b, a, y)
    freq_f, psd_f = signal.welch(yf, fs=fs, nperseg=512)

    # calculate filter response, create a results plot
    w, h = signal.freqz(b, a)
    wHz = w * fs / (2 * np.pi)
    fig, ax = open_figure('Notch Filter Example', 1, 2, figsize=(16, 6), sharex=True)
    ax[0].plot(wHz, 20 * np.log10(abs(h)), '-', lw=3)
    ax[1].semilogy(freq, psd, '.-', label='unfiltered')
    ax[1].semilogy(freq_f, psd_f, '.-', label='filtered')
    ax[1].legend(loc='upper left', bbox_to_anchor=(1, 1), shadow=True, numpoints=3)
    format_axes('freq, Hz', 'dB', 'Chebyshev II Bandstop Filter Response', ax[0])
    format_axes('freq, Hz', 'arb', axes=ax[1],
        title='Synthetic data\ntone at {}Hz should be filtered'.format(int(wf * fs / 2)))
    largefonts(16)
    fig.tight_layout()
    fig.subplots_adjust(right=0.8)
    plt.show()

def movingstd(x, k):
    """
    1d or 2d moving standard deviation in a rectangular window
    """

    from pandas import DataFrame
    from numpy import array

    if x.ndim == 1:
        return array(DataFrame(x).rolling(window=k, center=False).std()).squeeze()

    elif x.ndim == 2:
        return array(DataFrame(x).rolling(window=k, center=False).std())

    else:
        raise ValueError('too many dims')

def movingaverage(x, k):
    """
    1d or 2d moving average in a rectangular window

    for example,

    In [4]: a = np.array([0, 0, 1, 1, 1, 1, 0, 0])

    In [5]: b = a.reshape(-1, 2)

    In [6]: a
    Out[6]: array([0, 0, 1, 1, 1, 1, 0, 0])

    In [7]: b
    Out[7]:
    array([[0, 0],
           [1, 1],
           [1, 1],
           [0, 0]])

    In [8]: movingaverage(a, 3)
    Out[8]:
    array([        nan,         nan,  0.33333333,  0.66666667,  1.        ,
            1.        ,  0.66666667,  0.33333333])

    In [9]: movingaverage(a, 3).shape
    Out[9]: (8,)

    In [10]: movingaverage(b, 3)
    Out[10]:
    array([[        nan,         nan],
           [        nan,         nan],
           [ 0.66666667,  0.66666667],
           [ 0.66666667,  0.66666667]])

    In [11]: movingaverage(b, 3).shape
    Out[11]: (4, 2)
    """

    from pandas import DataFrame
    from numpy import array

    if x.ndim == 1:
        return array(DataFrame(x).rolling(window=k, center=False).mean()).squeeze()

    elif x.ndim == 2:
        return array(DataFrame(x).rolling(window=k, center=False).mean())

    else:
        raise ValueError('too many dims')

def movingsum(x, k):
    """
    1d or 2d moving sum in a rectangular window
    """

    from pandas import DataFrame
    from numpy import array

    if x.ndim == 1:
        return array(DataFrame(x).rolling(window=k, center=False).sum()).squeeze()

    elif x.ndim == 2:
        return array(DataFrame(x).rolling(window=k, center=False).sum())

    else:
        raise ValueError('too many dims')

def is_writeable(path):
    """
    basic approach to check if a path is writeable
    """
    import os

    try:
        file = os.path.join(path, r'delete_me.txt')
        fid = open(file, 'w')
        fid.close()
        os.remove(file)
        return True
    except PermissionError:
        return False

def return_dict(x):
    """
    recursive function that converts all defaultdicts to regular dicts
    this is useful as pickle sometimes has problems with defaultdicts
    """

    from collections import defaultdict
    for key in x.keys():
        if isinstance(x[key], defaultdict) or isinstance(x[key], dict):
            x[key] = return_dict(x[key])
    return dict(x)

def search_tree(datadir, pattern, use_pbar=False):
    """
    recursive search in datadir for pattern, with or without a progressbar
    """

    import os
    from glob import glob

    # 1 line solution with no progressbar
    if not use_pbar:
        return glob(os.path.join(datadir, '**', pattern), recursive=True)

    # otherwise must use os.walk to accomodate a progressbar
    from tqdm import tqdm
    from itertools import chain

    # get the toplevel directory list, initialize a 'found' list, counter, and progressbar
    toplevel = [os.path.join(datadir, x) for x in os.listdir(datadir)]
    toplevel = [x for x in toplevel if os.path.isdir(x)]
    found = []
    pbar = tqdm(desc='scanning datadir', total=len(toplevel))

    # walk thru data dir
    for root, _, _ in os.walk(datadir):

        # update the progressbar only for advances in toplevel directories, which could lead
        # to an inaccurate progressbar if the toplevel dirs do not have similar content
        if root in toplevel:
            pbar.update()

        # use glob to find search within root for pattern
        inside = glob(os.path.join(root, pattern))
        if inside:
            found.append(inside)

    # return a flattend list of lists
    return list(chain(*found))

def TIR(data):
    """
    return the total-inclusive range of data, works for 1d data without nans
    """

    import numpy as np

    return np.max(data) - np.min(data)

def get_utc_times(dt, zone, expected_zone=None):
    """
    given a datetime object 'dt' and a common time zone abbreviation 'zone' (e.g. PST, EST, CST, GMT, etc.),
    create a pandas DataFrame of all possible time differences with respect to UTC

    there is more than one possible time difference because every common time abbreviation can imply
    more than one unique timezone
    -- for example, PST can imply timezones of America/Los_Angeles or America/Juneau or many more!

    the datetime object must be provided (a single table from a single time-reference is not good enough)
    because daylight savings time makes this conversion dynamic with time

    the 'expected_zone' keyword argument limits the results to time zones that contain the 'expected_zone' string with unique UTC offset,
    this can greatly reduce the number or possible time zones returned
    (complicated but that is the only way with time zones)
    """

    import pytz
    import collections
    import numpy as np
    import pandas as pd
    from datetime import datetime

    # it is easy if zone is already recognized in pytz.all_timezones
    if zone in pytz.all_timezones:
        tzones = [zone]

    # otherwise get unambiguous time zone names associated with 'zone'
    else:

        # from https://stackoverflow.com/questions/36067621/python-all-possible-timezone-abbreviations-for-given-timezone-name-and-vise-ve
        # it is a mapping of how every official time zone in the world defines their time relative to their perceived definitions of standard time zones like PST, EST, ...
        # and also of how their time officially relates to the UTC international standard (same as GMT in practice)
        tzones = collections.defaultdict(set)
        for name in pytz.all_timezones:
            tzone = pytz.timezone(name)
            for utcoffset, dstoffset, tzabbrev in getattr(
                    tzone, '_transition_info', [[None, None, datetime.now(tzone).tzname()]]):
                tzones[tzabbrev].add(name)
        tzones = sorted(tzones[zone])

    # get the UTC time in each unambiguous time zone name
    utc = pytz.timezone('UTC')
    dts_utc = [pytz.timezone(x).localize(dt).astimezone(utc) for x in tzones]

    # find the hours of time difference relative to UTC in each unambiguous time zone name
    hrs = dt - np.array([x.replace(tzinfo=None) for x in dts_utc])
    hrs = np.array([int(x.total_seconds() / 3600) for x in hrs])

    # return results in a DataFrame
    results = pd.DataFrame(
                data=
                    {'datetime': np.tile(dt, hrs.size),
                    'zone': np.tile(zone, hrs.size),
                    'equivalent zones': tzones,
                    'datetime UTC': dts_utc,
                    'hrs wrt UTC': hrs},
                columns=['datetime', 'zone', 'equivalent zones', 'datetime UTC', 'hrs wrt UTC'])

    # adjust tzones based on the 'expected_zone' keyword argument
    if expected_zone is not None:

        # filter for similar names
        idx = results['equivalent zones'].str.lower().str.contains(expected_zone.lower())
        results = results[idx]

        # filter for unique offsets wrt UTC
        uhrs = pd.unique(results['hrs wrt UTC'])
        uresults = pd.DataFrame()
        equivalent_zones = []
        for uhr in uhrs:
            equivalent_zones.append(results[results['hrs wrt UTC'] == uhr]['equivalent zones'].iloc[0])
        idx = results['equivalent zones'].str.contains('|'.join(equivalent_zones))
        results = results[idx]

    return results

def get_unambiguous_tz(tz):
    """
    return an unambiguous timezone string given the ambiguous timezone string 'tz'
    'get_utc_times' can be used to help explore the possible unambiguous timezone strings

    this function is based on the author's interpretation of the intended meaning of ambiguous timezones
    the intended meaning of ambiguous timezones for other users can be different
    """

    import pytz
    from datetime import datetime

    # it is easy if tz is already recognized in pytz.all_timezones
    if tz in pytz.all_timezones:
        return tz

    if tz == 'PST':
        tz = 'US/Pacific'
    elif tz == 'JST':
        tz = 'Asia/Hong_Kong'
    elif tz == 'CST':
        tz = 'US/Central'
    elif tz == 'EDT':
        tz = 'America/New_York'
    elif tz == 'PDT':
        tz = 'America/Los_Angeles'
    elif tz == 'CEST':
        tz = 'Europe/Berlin'
    else:
        print('unknown and ambiguous timezone, more logic needed')
        return None

    assert get_utc_times(datetime.now(), tz)['equivalent zones'].str.contains(tz).any()
    return tz

def longest_substring_from_list(data):
    """
    return the longest common substring from a list of strings, inefficiently
    """

    def is_substr(find, data):
        """ used with longest_substring_from_list """
        if len(data) < 1 and len(find) < 1:
            return False
        for i in range(len(data)):
            if find not in data[i]:
                return False
        return True

    substr = ''
    if len(data) > 1 and len(data[0]) > 0:
        for i in range(len(data[0])):
            for j in range(len(data[0])-i+1):
                if j > len(substr) and is_substr(data[0][i:i+j], data):
                    substr = data[0][i:i+j]
    return substr

def memoized(f):
    """
    Memoization decorator for functions taking one or more arguments
    """
    class memodict(dict):
        def __init__(self, f):
            self.f = f
        def __call__(self, *args):
            return self[args]
        def __missing__(self, key):
            ret = self[key] = self.f(*key)
            return ret
    return memodict(f)

def pngs2ppt(pngs_dir, template=r'c:\pngs2ppt_template.pptx',
             fname=None, author='Author', date=None, title=None,
             add_exec_summary=False, img_width=9, takeaways=None, sort_order=None):

    """
    create powerpoint presentation from folder of png image files in pngs_dir
    supported keyword arguments:
    -- template, must point to the template file pngs2ppt_template.pptx
        get this file from the author if you do not have it
    -- fname, filename for output .pptx file, the default is to use the a filename of
        pngs2ppt.pptx in the same folder as pngs_dir
    -- author, name to include on the title page
    -- date, date to include on the title page, default is current date
    -- title, title for the title page, default is pngs2ppt
    -- add_exec_summary, True to add an 'Executive Summary' page to the presentation
    -- img_width, width of png images in the report, default is 9 inches
    -- takeaways, this is text added to each image slide, default is None
        if None, 'Takeaway Message' is added to each image slide
        if provided as a string, this same string is used for each slide
        if provided as a list of strings with length equal to number of png images,
        the corresponding string is added to each slide
    """

    import pptx
    from datetime import datetime
    from PIL import Image
    import os
    from glob import glob

    # set the fname, date, and title, if not provided as keyword arguments
    if fname is None:
        fname = os.path.join(pngs_dir, 'pngs2ppt.pptx')
    if date is None:
        date = datetime.now().strftime(r'%d %b %Y')
    if title is None:
        title = 'Created with python-pptx {}'.format(pptx.__version__)

    # get all png files in a list, return if empty
    pngs = glob(os.path.join(pngs_dir, r'*.png'))
    if not sort_order:
        pngs = sorted(pngs)
    else:
        assert len(sort_order) == len(pngs)
        pngs = [pngs[x] for x in sort_order]
    if len(pngs) == 0:
        print('No png images found, returning without creating a presentation')
        return

    # open a presentation object from the template file, save as fname before proceeding
    prs = pptx.Presentation(template)
    prs.save(fname)

    # add a title slide with title, author, date
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = '{}\n\n{}'.format(author, date)

    # add an executive summary if requested
    if add_exec_summary:
        exec_slide = prs.slides.add_slide(prs.slide_layouts[1])
        exec_slide.shapes.title.text = 'Executive Summary'
        exec_slide.placeholders[1].text = '...'
        exec_slide.placeholders[13].text = 'Takeaway'

    # get the width, height of a slide, get the native unit corresponding to 1 inch
    width, height = prs.slide_width, prs.slide_height
    inch = pptx.util.Inches(1)

    # process the takeaways keyword argument at this point
    if isinstance(takeaways, str):
        takeaways = [takeaways for _ in range(len(pngs))]
    elif isinstance(takeaways, list) and len(takeaways) == len(pngs):
        pass
    else:
        takeaways = [None for _ in range(len(pngs))]

    # add an image slide for each png file
    for png, takeaway in zip(pngs, takeaways):

        # get size of the png image
        png_width_nom, png_height_nom = Image.open(png).size

        # set the desired image width, and height to preserve the original aspect ratio
        png_width_actual = img_width * inch
        png_height_actual = png_height_nom * png_width_actual / png_width_nom

        # set the 'top' and 'left' keyword arguments to center the image
        top = int((height - png_height_actual) / 2)
        left = int((width - png_width_actual) / 2)

        # add the slide with the centered png image
        png_slide = prs.slides.add_slide(prs.slide_layouts[2])
        png_slide.shapes.title.text = os.path.split(png)[1][:-4]
        png_slide.shapes.add_picture(png, left=left, top=top, width=png_width_actual)

        # add the takeaway message
        if takeaway is None:
            tmp = png_slide.placeholders[13].element
            tmp.getparent().remove(tmp)
        else:
            png_slide.placeholders[13].text = takeaway

    # final save, print save location
    prs.save(fname)
    print('Presentation saved to {}'.format(fname))
